"""
RAG API路由
企业级RAG增强功能API
"""
from fastapi import APIRouter, Depends, HTTPException, Query, UploadFile, File, Header, Request, Path
from typing import Optional, List, Dict, Any
from datetime import datetime
from pydantic import BaseModel
import logging

from app.schemas.rag_schemas import (
    SearchRequest, SearchResponse, SearchChunk,
    IndexDocumentRequest, IndexBatchRequest, IndexResponse, IndexBatchResponse,
    DeleteDocumentRequest, DeleteResponse,
    CollectionStatsResponse,
    RAGConfigResponse, RAGHealthResponse,
    PatentSearchRequest, PatentSearchResponse, PatentDocumentResponse,
    PatentSearchAndIndexResponse, ConnectorInfo, DatabaseConnectorsResponse,
    AvailableModelsResponse, ModelOption,
    URLCrawlRequest, URLCrawlResponse,
    FileUploadResponse,
    ModelTestRequest, ModelTestResponse,
    CustomModelConfig, ModelConfigListResponse,
    PatentFullCrawlRequest, PatentIncrementalCrawlRequest, PatentCrawlStatusResponse,
    PatentSourceConfigCreate, PatentSourceConfigUpdate, PatentSourceConfigResponse, PatentSourceConfigListResponse
)
from app.schemas.rag_schemas import (
    SearchRequest, SearchResponse, SearchChunk,
    IndexDocumentRequest, IndexBatchRequest, IndexResponse, IndexBatchResponse,
    DeleteDocumentRequest, DeleteResponse,
    CollectionStatsResponse,
    RAGConfigResponse, RAGHealthResponse,
    PatentSearchRequest, PatentSearchResponse, PatentDocumentResponse,
    PatentSearchAndIndexResponse, ConnectorInfo, DatabaseConnectorsResponse,
    AvailableModelsResponse, ModelOption,
    URLCrawlRequest, URLCrawlResponse,
    FileUploadResponse,
    ModelTestRequest, ModelTestResponse,
    CustomModelConfig, ModelConfigListResponse,
    PatentFullCrawlRequest, PatentIncrementalCrawlRequest, PatentCrawlStatusResponse,
    PatentSourceConfigCreate, PatentSourceConfigUpdate, PatentSourceConfigResponse, PatentSourceConfigListResponse,
    # 新增的文件管理schemas
    RagDocumentResponse, RagDocumentListResponse, ReindexResponse
)
from app.database.models import RagDocument
from app.database.engine import async_session_factory
from app.ai.rag.enterprise_vector_db import enterprise_vector_db
from app.ai.rag.config import get_rag_settings, AVAILABLE_EMBEDDING_MODELS, AVAILABLE_RERANK_MODELS
from app.ai.rag.retrieval_pipeline import hybrid_search_engine

logger = logging.getLogger(__name__)

router = APIRouter(tags=["RAG增强"])


# ============== RAG搜索API ==============

@router.post("/search", response_model=SearchResponse)
async def search_documents(
    request: SearchRequest,
    tenant_id: Optional[str] = Header(None, alias="X-Tenant-ID", description="租户ID")
) -> SearchResponse:
    """
    RAG语义搜索
    
    - **query**: 搜索查询
    - **top_k**: 返回结果数量
    - **search_type**: 搜索类型 (semantic/keyword/hybrid)
    - **use_rerank**: 是否使用重排序
    - **filters**: 元数据过滤条件
    """
    from app.core.cache import cache
    
    try:
        # 尝试从缓存获取
        cache_key = f"rag:search:{tenant_id or 'default'}:{request.query}:{request.top_k}:{request.search_type}:{request.use_rerank}"
        cached_result = await cache.get(cache_key)
        if cached_result:
            import json
            data = json.loads(cached_result)
            return SearchResponse(**data)
        
        # 初始化混合搜索引擎
        hybrid_search_engine.initialize()
        
        # 执行搜索
        result = await hybrid_search_engine.search(
            query=request.query,
            top_k=request.top_k,
            tenant_id=tenant_id,
            search_type=request.search_type,
            use_rerank=request.use_rerank,
            filter_metadata=request.filters
        )
        
        # 转换为响应格式
        chunks = [
            SearchChunk(
                id=chunk.id,
                content=chunk.content,
                score=chunk.rerank_score or chunk.score,
                metadata=chunk.metadata,
                source=chunk.source
            )
            for chunk in result.chunks
        ]
        
        # 生成引用信息
        citations = None
        reference_list = None
        inline_citations = None
        
        if chunks:
            try:
                from app.ai.rag.citation_formatter import CitationFormatter, CitationManager
                
                # 创建引用管理器
                manager = CitationManager()
                
                for chunk in chunks:
                    # 为每个片段创建引用
                    manager.add_citation(
                        source_id=chunk.id,
                        source_type=chunk.metadata.get("source_type", "document"),
                        title=chunk.metadata.get("title", "Unknown"),
                        authors=chunk.metadata.get("authors"),
                        publication_date=chunk.metadata.get("publication_date"),
                        source=chunk.metadata.get("source"),
                        patent_number=chunk.metadata.get("patent_number"),
                        application_number=chunk.metadata.get("application_number"),
                        relevance_score=chunk.score
                    )
                
                # 生成引用格式
                citations = []
                for citation in manager.citations:
                    citations.append({
                        "index": manager.get_citation_index(citation.source_id),
                        "source_id": citation.source_id,
                        "title": citation.title,
                        "authors": citation.authors,
                        "publication_date": citation.publication_date,
                        "source": citation.source,
                        "relevance_score": citation.relevance_score
                    })
                
                reference_list = manager.format_reference_list(include_scores=True)
                inline_citations = CitationFormatter.create_inline_citations(manager.citations)
                
            except Exception as e:
                logger.warning(f"生成引用失败: {e}")
        
        response = SearchResponse(
            query=result.query,
            chunks=chunks,
            total=result.total,
            search_type=result.search_type,
            latency_ms=result.latency_ms,
            tenant_id=result.tenant_id,
            citations=citations,
            reference_list=reference_list,
            inline_citations=inline_citations
        )
        
        # 缓存结果 (5分钟TTL)
        import json
        await cache.set(cache_key, response.model_dump_json(), expire=300)

        # 保存搜索历史（异步执行，不阻塞响应）
        if tenant_id:
            try:
                from app.database.engine import async_session_factory
                from app.database.models import SearchHistory
                
                async with async_session_factory() as db:
                    history = SearchHistory(
                        tenant_id=int(tenant_id),
                        query=request.query,
                        search_type=request.search_type or "hybrid",
                        result_count=len(chunks),
                        filters=request.filters,
                        latency_ms=result.latency_ms
                    )
                    db.add(history)
                    await db.commit()
            except Exception as e:
                logger.warning(f"保存搜索历史失败: {e}")
        
        return response
        
    except Exception as e:
        logger.error(f"RAG搜索失败: {e}")
        raise HTTPException(status_code=500, detail=f"搜索失败: {str(e)}")


@router.post("/index", response_model=IndexResponse)
async def index_document(
    request: IndexDocumentRequest,
    tenant_id: Optional[str] = Header(None, alias="X-Tenant-ID", description="租户ID")
) -> IndexResponse:
    """
    索引单个专利文档到RAG系统
    
    - **document_id**: 文档ID
    - **application_number**: 申请号
    - **title**: 标题
    - **abstract**: 摘要
    - **claims**: 权利要求
    - **description**: 说明书
    - **technical_field**: 技术领域
    """
    try:
        # 构建文档
        document = {
            "id": request.document_id or request.application_number or "unknown",
            "application_number": request.application_number,
            "title": request.title,
            "abstract": request.abstract,
            "claims": request.claims,
            "description": request.description,
            "technical_field": request.technical_field,
            **(request.metadata or {})
        }
        
        # 索引文档
        success = await enterprise_vector_db.index_patent_document(
            document=document,
            tenant_id=tenant_id
        )
        
        return IndexResponse(
            success=success,
            document_id=document["id"],
            chunk_count=1,
            message="文档已索引" if success else "索引失败"
        )
        
    except Exception as e:
        logger.error(f"索引文档失败: {e}")
        raise HTTPException(status_code=500, detail=f"索引失败: {str(e)}")


@router.post("/index/batch", response_model=IndexBatchResponse)
async def index_documents_batch(
    request: IndexBatchRequest,
    tenant_id: Optional[str] = Header(None, alias="X-Tenant-ID", description="租户ID")
) -> IndexBatchResponse:
    """
    批量索引专利文档
    """
    try:
        documents = []
        for doc_req in request.documents:
            doc = {
                "id": doc_req.document_id or doc_req.application_number or "unknown",
                "application_number": doc_req.application_number,
                "title": doc_req.title,
                "abstract": doc_req.abstract,
                "claims": doc_req.claims,
                "description": doc_req.description,
                "technical_field": doc_req.technical_field,
                **(doc_req.metadata or {})
            }
            documents.append(doc)
        
        # 批量索引
        results = await enterprise_vector_db.index_patent_batch(
            documents=documents,
            tenant_id=tenant_id
        )
        
        successful = sum(1 for v in results.values() if v)
        
        return IndexBatchResponse(
            total=len(documents),
            successful=successful,
            failed=len(documents) - successful,
            results=results
        )
        
    except Exception as e:
        logger.error(f"批量索引失败: {e}")
        raise HTTPException(status_code=500, detail=f"批量索引失败: {str(e)}")


@router.delete("/documents", response_model=DeleteResponse)
async def delete_document(
    request: DeleteDocumentRequest,
    tenant_id: Optional[str] = Header(None, alias="X-Tenant-ID", description="租户ID")
) -> DeleteResponse:
    """
    删除索引文档
    """
    try:
        success = await enterprise_vector_db.delete_document(
            doc_id=request.document_id,
            tenant_id=tenant_id
        )
        
        return DeleteResponse(
            success=success,
            document_id=request.document_id,
            message="文档已删除" if success else "删除失败"
        )
        
    except Exception as e:
        logger.error(f"删除文档失败: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/stats", response_model=CollectionStatsResponse)
async def get_collection_stats(
    tenant_id: Optional[str] = Header(None, alias="X-Tenant-ID", description="租户ID")
) -> CollectionStatsResponse:
    """
    获取RAG集合统计信息
    """
    try:
        stats = await enterprise_vector_db.get_collection_stats(tenant_id=tenant_id)
        
        return CollectionStatsResponse(
            collection=stats.get("collection", "unknown"),
            document_count=stats.get("document_count", 0),
            backend=stats.get("backend", "chroma"),
            additional_info=stats
        )
        
    except Exception as e:
        logger.error(f"获取统计失败: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/config", response_model=RAGConfigResponse)
async def get_rag_config() -> RAGConfigResponse:
    """
    获取RAG配置信息
    """
    settings = get_rag_settings()
    
    return RAGConfigResponse(
        vector_db_type=settings.VECTOR_DB_TYPE,
        embedding_model=settings.EMBEDDING_MODEL,
        embedding_dimension=settings.EMBEDDING_DIMENSION,
        chunk_size=settings.CHUNK_SIZE,
        chunk_overlap=settings.CHUNK_OVERLAP,
        retrieval_top_k=settings.RETRIEVAL_TOP_K,
        rerank_enabled=settings.RERANK_ENABLED,
        hybrid_search_alpha=settings.HYBRID_SEARCH_ALPHA
    )


@router.put("/config", response_model=RAGConfigResponse)
async def update_rag_config(
    config: dict
) -> RAGConfigResponse:
    """
    更新RAG配置信息
    """
    from app.ai.rag.config import get_rag_settings
    
    rag_settings = get_rag_settings()
    
    # 更新配置
    if "vector_db_type" in config:
        rag_settings.VECTOR_DB_TYPE = config["vector_db_type"]
    if "embedding_model" in config:
        rag_settings.EMBEDDING_MODEL = config["embedding_model"]
    if "embedding_dimension" in config:
        rag_settings.EMBEDDING_DIMENSION = config["embedding_dimension"]
    if "chunk_size" in config:
        rag_settings.CHUNK_SIZE = config["chunk_size"]
    if "chunk_overlap" in config:
        rag_settings.CHUNK_OVERLAP = config["chunk_overlap"]
    if "retrieval_top_k" in config:
        rag_settings.RETRIEVAL_TOP_K = config["retrieval_top_k"]
    if "rerank_enabled" in config:
        rag_settings.RERANK_ENABLED = config["rerank_enabled"]
    if "hybrid_search_alpha" in config:
        rag_settings.HYBRID_SEARCH_ALPHA = config["hybrid_search_alpha"]
    
    # 保存配置到文件
    from app.ai.rag.config import save_rag_settings
    save_rag_settings()
    
    return RAGConfigResponse(
        vector_db_type=rag_settings.VECTOR_DB_TYPE,
        embedding_model=rag_settings.EMBEDDING_MODEL,
        embedding_dimension=rag_settings.EMBEDDING_DIMENSION,
        chunk_size=rag_settings.CHUNK_SIZE,
        chunk_overlap=rag_settings.CHUNK_OVERLAP,
        retrieval_top_k=rag_settings.RETRIEVAL_TOP_K,
        rerank_enabled=rag_settings.RERANK_ENABLED,
        hybrid_search_alpha=rag_settings.HYBRID_SEARCH_ALPHA
    )


@router.get("/models", response_model=AvailableModelsResponse)
async def get_available_models() -> AvailableModelsResponse:
    """
    获取可用的Embedding和Rerank模型列表
    """
    from app.ai.rag.config import AVAILABLE_EMBEDDING_MODELS, AVAILABLE_RERANK_MODELS
    
    return AvailableModelsResponse(
        embedding_models=[ModelOption(**m) for m in AVAILABLE_EMBEDDING_MODELS],
        rerank_models=[ModelOption(**m) for m in AVAILABLE_RERANK_MODELS]
    )


# ============== 自定义模型配置API ==============

@router.get("/models/custom", response_model=ModelConfigListResponse)
async def list_custom_model_configs() -> ModelConfigListResponse:
    """
    获取所有自定义模型配置
    """
    from app.ai.rag.config import list_custom_model_configs
    
    configs = await list_custom_model_configs()
    return ModelConfigListResponse(
        custom_configs=[CustomModelConfig(**c) for c in configs]
    )


@router.post("/models/custom", response_model=CustomModelConfig)
async def add_custom_model_config(
    config: CustomModelConfig
) -> CustomModelConfig:
    """
    添加自定义模型配置
    
    - **model_name**: 模型名称
    - **provider**: 提供商 (custom)
    - **api_url**: API URL (OpenAI兼容)
    - **api_key**: API Key
    - **dimension**: 向量维度
    - **enabled**: 是否启用
    """
    from app.ai.rag.config import set_custom_model_config
    
    config_dict = config.model_dump()
    await set_custom_model_config(config.model_name, config_dict)
    
    logger.info(f"添加自定义模型配置: {config.model_name}")
    return config


@router.put("/models/custom/{model_name}", response_model=CustomModelConfig)
async def update_custom_model_config(
    model_name: str,
    config: CustomModelConfig
) -> CustomModelConfig:
    """
    更新自定义模型配置
    """
    from app.ai.rag.config import set_custom_model_config
    
    config_dict = config.model_dump()
    await set_custom_model_config(model_name, config_dict)
    
    logger.info(f"更新自定义模型配置: {model_name}")
    return config


@router.delete("/models/custom/{model_name}")
async def delete_custom_model_config(model_name: str) -> dict:
    """
    删除自定义模型配置
    """
    from app.ai.rag.config import delete_custom_model_config
    
    await delete_custom_model_config(model_name)
    logger.info(f"删除自定义模型配置: {model_name}")
    return {"success": True, "message": f"已删除自定义模型配置: {model_name}"}


# ============== 自定义Reranker模型配置API ==============

class RerankerModelConfig(BaseModel):
    """Reranker模型配置"""
    model_name: str
    provider: str = "custom"
    api_url: Optional[str] = None
    api_key: Optional[str] = None
    enabled: bool = True
    description: Optional[str] = None


class RerankerConfigListResponse(BaseModel):
    """Reranker配置列表响应"""
    custom_configs: List[RerankerModelConfig] = []


@router.get("/rerankers/custom", response_model=RerankerConfigListResponse)
async def list_custom_reranker_configs() -> RerankerConfigListResponse:
    """获取所有自定义Reranker模型配置"""
    from app.ai.rag.config import list_custom_reranker_configs
    
    configs = await list_custom_reranker_configs()
    return RerankerConfigListResponse(
        custom_configs=[RerankerModelConfig(**c) for c in configs]
    )


@router.post("/rerankers/custom", response_model=RerankerModelConfig)
async def add_custom_reranker_config(config: RerankerModelConfig) -> RerankerModelConfig:
    """添加自定义Reranker模型配置"""
    from app.ai.rag.config import set_custom_reranker_config
    
    config_dict = config.model_dump()
    await set_custom_reranker_config(config.model_name, config_dict)
    
    logger.info(f"添加自定义Reranker模型配置: {config.model_name}")
    return config


@router.put("/rerankers/custom/{model_name}", response_model=RerankerModelConfig)
async def update_custom_reranker_config(
    model_name: str,
    config: RerankerModelConfig
) -> RerankerModelConfig:
    """更新自定义Reranker模型配置"""
    from app.ai.rag.config import set_custom_reranker_config
    
    config_dict = config.model_dump()
    await set_custom_reranker_config(model_name, config_dict)
    
    logger.info(f"更新自定义Reranker模型配置: {model_name}")
    return config


@router.delete("/rerankers/custom/{model_name}")
async def delete_custom_reranker_config(model_name: str) -> dict:
    """删除自定义Reranker模型配置"""
    from app.ai.rag.config import delete_custom_reranker_config
    
    await delete_custom_reranker_config(model_name)
    logger.info(f"删除自定义Reranker模型配置: {model_name}")
    return {"success": True, "message": f"已删除自定义Reranker模型配置: {model_name}"}


@router.get("/health", response_model=RAGHealthResponse)
async def rag_health_check() -> RAGHealthResponse:
    """
    RAG服务健康检查
    """
    from datetime import datetime
    
    # 获取配置
    settings = get_rag_settings()
    
    # 检查向量数据库连接状态
    vector_db_status = "healthy"
    try:
        if settings.VECTOR_DB_TYPE == "chroma":
            if settings.CHROMA_HOST:
                # 测试HTTP连接 - ChromaDB使用 /api/v2/heartbeat
                import httpx
                async with httpx.AsyncClient() as client:
                    try:
                        response = await client.get(
                            f"http://{settings.CHROMA_HOST}:{settings.CHROMA_PORT}/api/v2/heartbeat",
                            timeout=5.0
                        )
                        if response.status_code >= 400:
                            vector_db_status = "unhealthy"
                    except Exception as e:
                        # 连接失败
                        logger.warning(f"ChromaDB连接失败: {e}")
                        vector_db_status = "unhealthy"
            else:
                # 嵌入式模式，检查本地存储
                import os
                persist_dir = settings.get_chroma_persist_dir()
                if not os.path.exists(persist_dir):
                    vector_db_status = "unconfigured"
    except Exception as e:
        logger.warning(f"向量数据库健康检查失败: {e}")
        vector_db_status = "unhealthy"
    
    # 获取专利数据库连接器状态
    from app.ai.patent_database_api import patent_aggregator
    patent_aggregator.initialize()
    connectors = patent_aggregator.list_sources()
    
    return RAGHealthResponse(
        status=vector_db_status,
        vector_db=settings.VECTOR_DB_TYPE,
        embedding_service=settings.EMBEDDING_MODEL,
        connectors=[
            ConnectorInfo(name=c["name"], available=c["available"])
            for c in connectors
        ],
        timestamp=datetime.utcnow().isoformat()
    )


# ============== 专利数据库搜索API ==============

@router.post("/patents/search", response_model=PatentSearchResponse)
async def search_patents(
    request: PatentSearchRequest,
    tenant_id: Optional[str] = Header(None, alias="X-Tenant-ID", description="租户ID")
) -> PatentSearchResponse:
    """
    搜索公开专利数据库
    
    支持的数据源:
    - **dawei**: 大为专利数据库 (中国)
    - **cnipa**: CNIPA官方数据库 (中国)
    - **uspto**: USPTO PatentsView (美国)
    - **wipo**: WIPO PATENTSCOPE (国际)
    - **epo**: 欧洲专利局
    
    - **query**: 搜索查询
    - **sources**: 数据源列表
    - **max_results**: 每个数据源最大结果数
    - **date_from**: 起始日期 (YYYY-MM-DD)
    - **date_to**: 结束日期 (YYYY-MM-DD)
    - **classifications**: 专利分类
    - **auto_index**: 是否自动索引到RAG系统
    """
    try:
        from app.ai.patent_database_api import patent_aggregator
        
        # 执行搜索 (cast to List[str] to satisfy type checker)
        sources_arg: Optional[List[str]] = request.sources  # type: ignore
        results = await patent_aggregator.search_all(
            query=request.query,
            sources=sources_arg,
            max_results_per_source=request.max_results
        )
        
        # 转换为响应格式
        response_results = {}
        total_count = 0
        
        for source, patents in results.items():
            response_results[source] = [
                PatentDocumentResponse(
                    id=p.publication_number or p.application_number or "",
                    application_number=p.application_number,
                    publication_number=p.publication_number,
                    title=p.title,
                    abstract=p.abstract,
                    claims=p.claims.split("\n") if p.claims else None,
                    applicant=p.applicant,
                    assignee=p.assignee,
                    inventor=p.inventor,
                    application_date=p.application_date,
                    publication_date=p.publication_date,
                    ipc_classification=p.ipc_classification.split(";") if p.ipc_classification else None,
                    cpc_classification=p.cpc_classification.split(";") if p.cpc_classification else None,
                    patent_type=p.patent_type,
                    source=p.source,
                    url=p.url
                )
                for p in patents
            ]
            total_count += len(patents)
        
        # 如果需要自动索引
        if request.auto_index and tenant_id:
            from app.ai.rag.enterprise_vector_db import enterprise_vector_db
            
            for source, patents in results.items():
                if patents:
                    docs = [p.to_dict() for p in patents]
                    await enterprise_vector_db.index_patent_batch(docs, tenant_id=tenant_id)
        
        return PatentSearchResponse(
            query=request.query,
            results=response_results,
            total_count=total_count,
            sources_searched=list(results.keys())
        )
        
    except Exception as e:
        logger.error(f"专利搜索失败: {e}")
        raise HTTPException(status_code=500, detail=f"搜索失败: {str(e)}")


@router.get("/patents/sources", response_model=DatabaseConnectorsResponse)
async def list_patent_sources(
    tenant_id: str = Query(None, description="租户ID")
) -> DatabaseConnectorsResponse:
    """
    列出所有可用的专利数据源
    """
    from app.ai.patent_database_api import patent_aggregator
    patent_aggregator.initialize()
    
    # 加载用户配置的数据源
    if tenant_id:
        try:
            patent_aggregator.load_user_configs(int(tenant_id))
        except Exception as e:
            pass  # 忽略配置加载错误
    
    sources = patent_aggregator.list_sources()
    
    return DatabaseConnectorsResponse(
        connectors=[
            ConnectorInfo(name=s["name"], available=s["available"])
            for s in sources
        ]
    )


# ─── 专利数据源配置 CRUD ─────────────────────────────


@router.get("/patents/source-configs", response_model=PatentSourceConfigListResponse)
async def list_patent_source_configs(
    tenant_id: str = Query(..., description="租户ID")
) -> PatentSourceConfigListResponse:
    """
    获取专利数据源配置列表
    """
    from sqlalchemy import select
    from app.database.models import PatentSourceConfig
    from app.database.engine import async_session_factory
    
    async with async_session_factory() as db:
        # 查询该租户的配置
        result = await db.execute(
            select(PatentSourceConfig).where(
                PatentSourceConfig.tenant_id == int(tenant_id) if tenant_id else True
            )
        )
        configs = result.scalars().all()
        
        return PatentSourceConfigListResponse(
            configs=[PatentSourceConfigResponse.model_validate(c) for c in configs],
            total=len(configs)
        )


@router.post("/patents/source-configs", response_model=PatentSourceConfigResponse)
async def create_patent_source_config(
    config: PatentSourceConfigCreate,
    tenant_id: str = Query(..., description="租户ID")
) -> PatentSourceConfigResponse:
    """
    创建专利数据源配置
    """
    from sqlalchemy import select
    from app.database.models import PatentSourceConfig
    from app.database.engine import async_session_factory
    
    async with async_session_factory() as db:
        # 检查是否已存在相同数据源配置
        result = await db.execute(
            select(PatentSourceConfig).where(
                PatentSourceConfig.source_name == config.source_name,
                PatentSourceConfig.tenant_id == int(tenant_id) if tenant_id else True
            )
        )
        existing = result.scalar_one_or_none()
        
        if existing:
            # 更新现有配置
            existing.api_key = config.api_key
            existing.app_id = config.app_id
            existing.api_url = config.api_url
            existing.is_enabled = config.is_enabled
            await db.commit()
            await db.refresh(existing)
            return PatentSourceConfigResponse.model_validate(existing)
        
        # 创建新配置
        new_config = PatentSourceConfig(
            tenant_id=int(tenant_id) if tenant_id else None,
            source_name=config.source_name,
            api_key=config.api_key,
            app_id=config.app_id,
            api_url=config.api_url,
            is_enabled=config.is_enabled
        )
        db.add(new_config)
        await db.commit()
        await db.refresh(new_config)
        
        return PatentSourceConfigResponse.model_validate(new_config)


@router.put("/patents/source-configs/{config_id}", response_model=PatentSourceConfigResponse)
async def update_patent_source_config(
    config_id: int,
    config: PatentSourceConfigUpdate,
    tenant_id: str = Query(..., description="租户ID")
) -> PatentSourceConfigResponse:
    """
    更新专利数据源配置
    """
    from sqlalchemy import select
    from app.database.models import PatentSourceConfig
    from app.database.engine import async_session_factory
    
    async with async_session_factory() as db:
        result = await db.execute(
            select(PatentSourceConfig).where(
                PatentSourceConfig.id == config_id,
                PatentSourceConfig.tenant_id == int(tenant_id) if tenant_id else True
            )
        )
        db_config = result.scalar_one_or_none()
        
        if not db_config:
            raise HTTPException(status_code=404, detail="配置不存在")
        
        # 更新非空字段
        if config.api_key is not None:
            db_config.api_key = config.api_key
        if config.app_id is not None:
            db_config.app_id = config.app_id
        if config.api_url is not None:
            db_config.api_url = config.api_url
        if config.is_enabled is not None:
            db_config.is_enabled = config.is_enabled
        
        await db.commit()
        await db.refresh(db_config)
        
        return PatentSourceConfigResponse.model_validate(db_config)


@router.delete("/patents/source-configs/{config_id}")
async def delete_patent_source_config(
    config_id: int,
    tenant_id: str = Query(..., description="租户ID")
):
    """
    删除专利数据源配置
    """
    from sqlalchemy import select, delete
    from app.database.models import PatentSourceConfig
    from app.database.engine import async_session_factory
    
    async with async_session_factory() as db:
        result = await db.execute(
            select(PatentSourceConfig).where(
                PatentSourceConfig.id == config_id,
                PatentSourceConfig.tenant_id == int(tenant_id) if tenant_id else True
            )
        )
        db_config = result.scalar_one_or_none()
        
        if not db_config:
            raise HTTPException(status_code=404, detail="配置不存在")
        
        await db.delete(db_config)
        await db.commit()
        
        return {"success": True, "message": "配置已删除"}


@router.post("/patents/search-and-index", response_model=PatentSearchAndIndexResponse)
async def search_and_index_patents(
    request: PatentSearchRequest,
    tenant_id: str = Query(..., description="租户ID")
) -> PatentSearchAndIndexResponse:
    """
    搜索专利并自动索引到RAG系统
    """
    try:
        from app.ai.patent_database_api import patent_aggregator
        
        # 搜索 (cast to List[str] to satisfy type checker)
        sources_arg: Optional[List[str]] = request.sources  # type: ignore
        results = await patent_aggregator.search_all(
            query=request.query,
            sources=sources_arg,
            max_results_per_source=request.max_results
        )
        
        # 收集所有专利
        all_patents = []
        for source, patents in results.items():
            all_patents.extend(patents)
        
        # 索引到向量数据库
        indexed_count = 0
        if all_patents:
            docs = [p.to_dict() for p in all_patents]
            index_results = await enterprise_vector_db.index_patent_batch(
                documents=docs,
                tenant_id=tenant_id
            )
            indexed_count = sum(1 for v in index_results.values() if v)
        
        return PatentSearchAndIndexResponse(
            total_found=len(all_patents),
            indexed=indexed_count,
            sources={k: len(v) for k, v in results.items()}
        )
        
    except Exception as e:
        logger.error(f"搜索并索引失败: {e}")
        raise HTTPException(status_code=500, detail=str(e))


# ============== 专利全量爬取API ==============

@router.post("/patents/crawl/full", response_model=PatentCrawlStatusResponse)
async def full_crawl_patents(
    request: PatentFullCrawlRequest
) -> PatentCrawlStatusResponse:
    """
    专利全量爬取 - 一次性大规模数据导入
    
    根据关键词搜索专利数据并批量索引到RAG系统
    
    - **query**: 搜索关键词 (如: 人工智能, 区块链, 新能源)
    - **sources**: 数据源列表
    - **max_results**: 每个数据源最大结果数 (1-1000)
    - **batch_size**: 批量索引大小 (10-100)
    - **tenant_id**: 租户ID
    """
    from datetime import datetime
    import asyncio
    
    start_time = datetime.now()
    total_found = 0
    indexed = 0
    failed = 0
    errors = []
    sources_count = {}
    
    try:
        from app.ai.patent_database_api import patent_aggregator
        
        logger.info(f"开始全量爬取: query={request.query}, sources={request.sources}")
        
        # 搜索专利
        results = await patent_aggregator.search_all(
            query=request.query,
            sources=request.sources,
            max_results_per_source=request.max_results
        )
        
        # 收集所有专利
        all_patents = []
        for source, patents in results.items():
            sources_count[source] = len(patents)
            total_found += len(patents)
            all_patents.extend(patents)
        
        logger.info(f"全量爬取找到 {total_found} 条专利")
        
        # 分批索引
        if all_patents:
            docs = [p.to_dict() for p in all_patents]
            
            # 分批处理
            for i in range(0, len(docs), request.batch_size):
                batch = docs[i:i + request.batch_size]
                try:
                    index_results = await enterprise_vector_db.index_patent_batch(
                        documents=batch,
                        tenant_id=request.tenant_id
                    )
                    indexed += sum(1 for v in index_results.values() if v)
                    failed += sum(1 for v in index_results.values() if not v)
                    
                    # 添加小延迟避免过载
                    await asyncio.sleep(0.5)
                    
                except Exception as batch_err:
                    logger.error(f"批次索引失败: {batch_err}")
                    failed += len(batch)
                    errors.append({"batch": i, "error": str(batch_err)})
        
        end_time = datetime.now()
        duration = (end_time - start_time).total_seconds()
        
        return PatentCrawlStatusResponse(
            status="completed",
            total_found=total_found,
            indexed=indexed,
            failed=failed,
            sources=sources_count,
            message=f"全量爬取完成，共找到{total_found}条专利，成功索引{indexed}条，耗时{duration:.1f}秒",
            started_at=start_time.isoformat(),
            completed_at=end_time.isoformat(),
            errors=errors if errors else None
        )
        
    except Exception as e:
        logger.error(f"全量爬取失败: {e}")
        return PatentCrawlStatusResponse(
            status="failed",
            total_found=total_found,
            indexed=indexed,
            failed=failed,
            sources=sources_count,
            message=f"爬取失败: {str(e)}",
            started_at=start_time.isoformat(),
            errors=[{"error": str(e)}]
        )


# ============== 专利增量爬取API ==============

@router.post("/patents/crawl/incremental", response_model=PatentCrawlStatusResponse)
async def incremental_crawl_patents(
    request: PatentIncrementalCrawlRequest
) -> PatentCrawlStatusResponse:
    """
    专利增量爬取 - 实时增量更新
    
    根据日期范围搜索新专利并增量索引
    
    - **query**: 搜索关键词
    - **sources**: 数据源列表
    - **date_from**: 起始日期 (YYYY-MM-DD), 默认最近7天
    - **max_results**: 每个数据源最大结果数 (1-500)
    - **tenant_id**: 租户ID
    """
    from datetime import datetime, timedelta
    import asyncio
    
    start_time = datetime.now()
    total_found = 0
    indexed = 0
    failed = 0
    errors = []
    sources_count = {}
    
    # 计算默认日期范围 (最近7天)
    if not request.date_from:
        date_from = (datetime.now() - timedelta(days=7)).strftime("%Y-%m-%d")
    else:
        date_from = request.date_from
    
    try:
        from app.ai.patent_database_api import patent_aggregator
        
        logger.info(f"开始增量爬取: query={request.query}, date_from={date_from}, sources={request.sources}")
        
        # 搜索专利 (带日期过滤)
        results = await patent_aggregator.search_all(
            query=request.query,
            sources=request.sources,
            max_results_per_source=request.max_results
        )
        
        # 过滤新专利 (日期大于等于date_from)
        filtered_patents = []
        for source, patents in results.items():
            source_patents = []
            for p in patents:
                # 检查发布日期
                pub_date = p.publication_date
                if pub_date:
                    try:
                        pub_dt = datetime.strptime(pub_date, "%Y-%m-%d")
                        from_dt = datetime.strptime(date_from, "%Y-%m-%d")
                        if pub_dt >= from_dt:
                            source_patents.append(p)
                    except:
                        # 无法解析日期, 默认包含
                        source_patents.append(p)
                else:
                    # 没有发布日期, 默认包含
                    source_patents.append(p)
            
            sources_count[source] = len(source_patents)
            total_found += len(source_patents)
            filtered_patents.extend(source_patents)
        
        logger.info(f"增量爬取找到 {total_found} 条新专利 (从 {date_from} 起)")
        
        # 索引新专利
        if filtered_patents:
            docs = [p.to_dict() for p in filtered_patents]
            
            for i in range(0, len(docs), 50):  # 默认批次大小50
                batch = docs[i:i + 50]
                try:
                    index_results = await enterprise_vector_db.index_patent_batch(
                        documents=batch,
                        tenant_id=request.tenant_id
                    )
                    indexed += sum(1 for v in index_results.values() if v)
                    failed += sum(1 for v in index_results.values() if not v)
                    
                    await asyncio.sleep(0.3)  # 较短延迟
                    
                except Exception as batch_err:
                    logger.error(f"批次索引失败: {batch_err}")
                    failed += len(batch)
                    errors.append({"batch": i, "error": str(batch_err)})
        
        end_time = datetime.now()
        duration = (end_time - start_time).total_seconds()
        
        return PatentCrawlStatusResponse(
            status="completed",
            total_found=total_found,
            indexed=indexed,
            failed=failed,
            sources=sources_count,
            message=f"增量爬取完成，从{date_from}起找到{total_found}条新专利，成功索引{indexed}条，耗时{duration:.1f}秒",
            started_at=start_time.isoformat(),
            completed_at=end_time.isoformat(),
            errors=errors if errors else None
        )
        
    except Exception as e:
        logger.error(f"增量爬取失败: {e}")
        return PatentCrawlStatusResponse(
            status="failed",
            total_found=total_found,
            indexed=indexed,
            failed=failed,
            sources=sources_count,
            message=f"爬取失败: {str(e)}",
            started_at=start_time.isoformat(),
            errors=[{"error": str(e)}]
        )


# ============== 文件上传API ==============

@router.post("/upload", response_model=FileUploadResponse)
async def upload_file(
    tenant_id: Optional[str] = Header(None, alias="X-Tenant-ID", description="租户ID")
) -> FileUploadResponse:
    """
    上传文件并提取文本内容
    
    支持的文件类型:
    - PDF (.pdf)
    - Word (.docx, .doc)
    - Text (.txt)
    - Excel (.xlsx, .xls)
    - PowerPoint (.pptx)
    - Images (.png, .jpg, .jpeg) - 需要OCR
    """
    from fastapi import UploadFile, File as FastAPIFile
    from fastapi.responses import JSONResponse
    
    # 接收文件
    files = await FastAPIFile().__class__(
        description="要上传的文件"
    ).read()
    
    # 注意: FastAPI文件上传需要特殊处理，这里需要用UploadFile
    pass


@router.post("/upload/text", response_model=FileUploadResponse)
async def upload_text_file(
    file: UploadFile = File(..., description="上传的文件"),
    tenant_id: Optional[str] = Header(None, alias="X-Tenant-ID", description="租户ID")
) -> FileUploadResponse:
    """
    上传文本文件并提取内容
    """
    import os
    import tempfile
    
    file_name = file.filename
    file_ext = os.path.splitext(file_name)[1].lower()
    
    # 保存上传的文件到临时目录
    with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as tmp:
        content = await file.read()
        tmp.write(content)
        tmp_path = tmp.name
    
    try:
        extracted_text = ""
        chunk_count = 0
        
        # 根据文件类型提取文本
        if file_ext == '.txt':
            with open(tmp_path, 'r', encoding='utf-8') as f:
                extracted_text = f.read()
            chunk_count = 1
            
        elif file_ext == '.pdf':
            try:
                from pypdf import PdfReader
                reader = PdfReader(tmp_path)
                for page in reader.pages:
                    extracted_text += page.extract_text() + "\n"
                chunk_count = len(reader.pages)
            except Exception as e:
                logger.warning(f"PDF提取失败: {e}")
                extracted_text = f"[PDF文件: {file_name}]"
                chunk_count = 1
                
        elif file_ext in ['.docx', '.doc']:
            try:
                from docx import Document
                doc = Document(tmp_path)
                for para in doc.paragraphs:
                    extracted_text += para.text + "\n"
                chunk_count = len(doc.paragraphs)
            except Exception as e:
                logger.warning(f"Word提取失败: {e}")
                extracted_text = f"[Word文档: {file_name}]"
                chunk_count = 1
                
        elif file_ext in ['.xlsx', '.xls']:
            try:
                import openpyxl
                wb = openpyxl.load_workbook(tmp_path)
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    for row in ws.iter_rows(values_only=True):
                        extracted_text += " | ".join([str(cell) if cell else "" for cell in row]) + "\n"
                chunk_count = wb.sheetnames.__len__() if hasattr(wb, 'sheetnames') else 1
            except Exception as e:
                logger.warning(f"Excel提取失败: {e}")
                extracted_text = f"[Excel文件: {file_name}]"
                chunk_count = 1
                
        else:
            extracted_text = f"[不支持的文件类型: {file_ext}]"
            chunk_count = 1
        
        return FileUploadResponse(
            success=True,
            file_name=file_name,
            file_type=file_ext,
            file_size=len(content),
            content=extracted_text[:10000] if extracted_text else None,  # 限制返回内容长度
            chunk_count=chunk_count
        )
        
    except Exception as e:
        logger.error(f"文件处理失败: {e}")
        return FileUploadResponse(
            success=False,
            file_name=file_name,
            file_type=file_ext,
            file_size=len(content) if 'content' in dir() else 0,
            error=str(e)
        )
    finally:
        # 清理临时文件
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)


@router.post("/upload/file-and-index")
async def upload_and_index_file(
    file: UploadFile = File(..., description="上传的文件"),
    tenant_id: str = Query(..., description="租户ID")
):
    """
    上传文件并自动索引到RAG系统
    
    - 保存文件到永久存储
    - 提取文本内容
    - 创建数据库记录
    - 索引到向量数据库
    """
    import os
    import tempfile
    import uuid
    from sqlalchemy import select
    
    file_name = file.filename
    file_ext = os.path.splitext(file_name)[1].lower()
    
    # 检查文件类型
    allowed_extensions = ['.txt', '.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx']
    if file_ext not in allowed_extensions:
        raise HTTPException(status_code=400, detail=f"不支持的文件类型: {file_ext}")
    
    # 创建永久存储目录
    storage_dir = os.path.join("uploads", "rag_files", str(tenant_id))
    os.makedirs(storage_dir, exist_ok=True)
    
    # 生成唯一文件名
    unique_filename = f"{uuid.uuid4().hex}_{file_name}"
    permanent_path = os.path.join(storage_dir, unique_filename)
    
    # 先保存到临时文件
    with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as tmp:
        content = await file.read()
        tmp.write(content)
        tmp_path = tmp.name
    
    try:
        extracted_text = ""
        chunk_count = 0
        
        # 根据文件类型提取文本
        if file_ext == '.txt':
            with open(tmp_path, 'r', encoding='utf-8') as f:
                extracted_text = f.read()
            chunk_count = 1
            
        elif file_ext == '.pdf':
            try:
                from pypdf import PdfReader
                reader = PdfReader(tmp_path)
                for page in reader.pages:
                    extracted_text += page.extract_text() + "\n"
                chunk_count = len(reader.pages)
            except Exception as e:
                logger.warning(f"PDF提取失败: {e}")
                extracted_text = f"[PDF文件: {file_name}]"
                chunk_count = 1
                
        elif file_ext in ['.docx', '.doc']:
            try:
                from docx import Document
                doc = Document(tmp_path)
                for para in doc.paragraphs:
                    extracted_text += para.text + "\n"
                chunk_count = len(doc.paragraphs) or 1
            except Exception as e:
                logger.warning(f"Word提取失败: {e}")
                extracted_text = f"[Word文档: {file_name}]"
                chunk_count = 1
                
        elif file_ext in ['.xlsx', '.xls']:
            try:
                import openpyxl
                wb = openpyxl.load_workbook(tmp_path)
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    for row in ws.iter_rows(values_only=True):
                        extracted_text += " | ".join([str(cell) if cell else "" for cell in row]) + "\n"
                chunk_count = len(wb.sheetnames) or 1
            except Exception as e:
                logger.warning(f"Excel提取失败: {e}")
                extracted_text = f"[Excel文件: {file_name}]"
                chunk_count = 1
        elif file_ext == '.pptx':
            try:
                from pptx import Presentation
                prs = Presentation(tmp_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                extracted_text += para.text + "\n"
                chunk_count = len(prs.slides) or 1
            except Exception as e:
                logger.warning(f"PPT提取失败: {e}")
                extracted_text = f"[PowerPoint文件: {file_name}]"
                chunk_count = 1
        else:
            extracted_text = f"[不支持的文件类型: {file_ext}]"
            chunk_count = 1
        
        # 将临时文件移动到永久存储
        import shutil
        shutil.move(tmp_path, permanent_path)
        tmp_path = None  # 标记已移动，不再清理
        
        # 创建数据库记录
        async with async_session_factory() as db:
            # 创建RagDocument记录
            rag_doc = RagDocument(
                tenant_id=int(tenant_id),
                file_name=file_name,
                file_type=file_ext,
                file_size=len(content),
                file_path=permanent_path,
                extracted_text=extracted_text[:50000] if extracted_text else None,  # 限制存储的文本长度
                chunk_count=chunk_count,
                indexed_count=0,
                status="indexing"
            )
            db.add(rag_doc)
            await db.commit()
            await db.refresh(rag_doc)
            
            document_id = rag_doc.id
        
        # 索引到向量数据库
        doc_id = f"file_{document_id}"
        document = {
            "id": doc_id,
            "title": file_name,
            "description": extracted_text,
            "metadata": {
                "source_type": "file_upload",
                "file_name": file_name,
                "file_type": file_ext,
                "document_id": document_id
            }
        }
        
        success = await enterprise_vector_db.index_patent_document(
            document=document,
            tenant_id=tenant_id
        )
        
        # 更新数据库记录状态
        async with async_session_factory() as db:
            result = await db.execute(
                select(RagDocument).where(RagDocument.id == document_id)
            )
            doc = result.scalar_one_or_none()
            if doc:
                if success:
                    doc.status = "completed"
                    doc.indexed_count = 1
                    doc.indexed_at = datetime.utcnow()
                else:
                    doc.status = "failed"
                    doc.error_message = "索引失败"
                await db.commit()
        
        return {
            "success": True,
            "document_id": document_id,
            "file_name": file_name,
            "file_type": file_ext,
            "file_size": len(content),
            "chunk_count": chunk_count,
            "message": "文件已上传并索引" if success else "索引失败"
        }
        
    except Exception as e:
        logger.error(f"文件上传并索引失败: {e}")
        # 清理永久文件
        if os.path.exists(permanent_path):
            try:
                os.unlink(permanent_path)
            except:
                pass
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        # 清理临时文件（如果还没移动）
        if tmp_path and os.path.exists(tmp_path):
            try:
                os.unlink(tmp_path)
            except:
                pass


# ============== 文件管理API ==============

@router.get("/files", response_model=RagDocumentListResponse)
async def list_rag_documents(
    tenant_id: str = Query(..., description="租户ID"),
    skip: int = Query(0, ge=0, description="跳过记录数"),
    limit: int = Query(50, ge=1, le=100, description="返回记录数"),
    status: Optional[str] = Query(None, description="过滤状态: pending, indexing, completed, failed")
) -> RagDocumentListResponse:
    """
    获取上传文件列表
    
    - **tenant_id**: 租户ID
    - **skip**: 跳过记录数
    - **limit**: 返回记录数 (最大100)
    - **status**: 按状态过滤
    """
    # 处理tenant_id转换
    try:
        tenant_id_int = int(tenant_id) if tenant_id != "default" else 1
    except ValueError:
        tenant_id_int = 1  # 默认租户
    from sqlalchemy import select, func, desc
    
    async with async_session_factory() as db:
        # 构建查询
        query = select(RagDocument).where(RagDocument.tenant_id == tenant_id_int)
        
        # 添加状态过滤
        if status:
            query = query.where(RagDocument.status == status)
        
        # 获取总数
        count_query = select(func.count()).select_from(RagDocument).where(
            RagDocument.tenant_id == int(tenant_id)
        )
        if status:
            count_query = count_query.where(RagDocument.status == status)
        
        total_result = await db.execute(count_query)
        total = total_result.scalar() or 0
        
        # 获取分页数据
        query = query.order_by(desc(RagDocument.created_at)).offset(skip).limit(limit)
        result = await db.execute(query)
        documents = result.scalars().all()
        
        return RagDocumentListResponse(
            documents=[
                RagDocumentResponse(
                    id=doc.id,
                    tenant_id=doc.tenant_id,
                    file_name=doc.file_name,
                    file_type=doc.file_type,
                    file_size=doc.file_size,
                    file_path=doc.file_path,
                    chunk_count=doc.chunk_count,
                    indexed_count=doc.indexed_count,
                    status=doc.status,
                    error_message=doc.error_message,
                    indexed_at=doc.indexed_at.isoformat() if doc.indexed_at else None,
                    created_at=doc.created_at.isoformat(),
                    updated_at=doc.updated_at.isoformat() if doc.updated_at else None
                )
                for doc in documents
            ],
            total=total
        )


@router.get("/files/{file_id}", response_model=RagDocumentResponse)
async def get_rag_document(
    file_id: int,
    tenant_id: str = Query(..., description="租户ID")
) -> RagDocumentResponse:
    """
    获取文件详情
    
    - **file_id**: 文件ID
    - **tenant_id**: 租户ID
    """
    from sqlalchemy import select
    
    async with async_session_factory() as db:
        result = await db.execute(
            select(RagDocument).where(
                RagDocument.id == file_id,
                RagDocument.tenant_id == int(tenant_id)
            )
        )
        doc = result.scalar_one_or_none()
        
        if not doc:
            raise HTTPException(status_code=404, detail="文件不存在")
        
        return RagDocumentResponse(
            id=doc.id,
            tenant_id=doc.tenant_id,
            file_name=doc.file_name,
            file_type=doc.file_type,
            file_size=doc.file_size,
            file_path=doc.file_path,
            chunk_count=doc.chunk_count,
            indexed_count=doc.indexed_count,
            status=doc.status,
            error_message=doc.error_message,
            indexed_at=doc.indexed_at.isoformat() if doc.indexed_at else None,
            created_at=doc.created_at.isoformat(),
            updated_at=doc.updated_at.isoformat() if doc.updated_at else None
        )


@router.delete("/files/{file_id}")
async def delete_rag_document(
    file_id: int,
    tenant_id: str = Query(..., description="租户ID")
) -> dict:
    """
    删除上传的文件
    
    - **file_id**: 文件ID
    - **tenant_id**: 租户ID
    
    删除向量数据库中的索引和物理文件
    """
    from sqlalchemy import select, delete
    import os
    
    async with async_session_factory() as db:
        # 查找文件记录
        result = await db.execute(
            select(RagDocument).where(
                RagDocument.id == file_id,
                RagDocument.tenant_id == int(tenant_id)
            )
        )
        doc = result.scalar_one_or_none()
        
        if not doc:
            raise HTTPException(status_code=404, detail="文件不存在")
        
        # 1. 从向量数据库删除索引
        try:
            doc_id = f"file_{doc.id}"
            await enterprise_vector_db.delete_document(
                doc_id=doc_id,
                tenant_id=tenant_id
            )
        except Exception as e:
            logger.warning(f"删除向量索引失败: {e}")
        
        # 2. 删除物理文件
        try:
            if doc.file_path and os.path.exists(doc.file_path):
                os.unlink(doc.file_path)
        except Exception as e:
            logger.warning(f"删除物理文件失败: {e}")
        
        # 3. 删除数据库记录
        await db.delete(doc)
        await db.commit()
        
        return {"success": True, "message": "文件已删除"}


@router.post("/files/{file_id}/reindex", response_model=ReindexResponse)
async def reindex_rag_document(
    file_id: int,
    tenant_id: str = Query(..., description="租户ID")
) -> ReindexResponse:
    """
    重新索引文件
    
    - **file_id**: 文件ID
    - **tenant_id**: 租户ID
    """
    from sqlalchemy import select
    import os
    
    async with async_session_factory() as db:
        result = await db.execute(
            select(RagDocument).where(
                RagDocument.id == file_id,
                RagDocument.tenant_id == int(tenant_id)
            )
        )
        doc = result.scalar_one_or_none()
        
        if not doc:
            raise HTTPException(status_code=404, detail="文件不存在")
        
        if not doc.extracted_text:
            raise HTTPException(status_code=400, detail="文件没有提取文本，无法重新索引")
        
        # 更新状态为 indexing
        doc.status = "indexing"
        await db.commit()
        
        try:
            # 重新索引文档
            doc_id = f"file_{doc.id}"
            document = {
                "id": doc_id,
                "title": doc.file_name,
                "description": doc.extracted_text,
                "metadata": {
                    "source_type": "file_upload",
                    "file_name": doc.file_name,
                    "file_type": doc.file_type
                }
            }
            
            success = await enterprise_vector_db.index_patent_document(
                document=document,
                tenant_id=tenant_id
            )
            
            if success:
                doc.status = "completed"
                doc.indexed_count = 1
                from datetime import datetime
                doc.indexed_at = datetime.utcnow()
            else:
                doc.status = "failed"
                doc.error_message = "索引失败"
            
            await db.commit()
            
            return ReindexResponse(
                success=success,
                document_id=doc.id,
                indexed_count=doc.indexed_count,
                message="文件重新索引成功" if success else "索引失败"
            )
            
        except Exception as e:
            doc.status = "failed"
            doc.error_message = str(e)
            await db.commit()
            raise HTTPException(status_code=500, detail=f"重新索引失败: {str(e)}")


@router.get("/files/{file_id}/download")
async def download_rag_document(
    file_id: int,
    tenant_id: str = Query(..., description="租户ID")
):
    """
    下载原始文件
    
    - **file_id**: 文件ID
    - **tenant_id**: 租户ID
    """
    from sqlalchemy import select
    from fastapi.responses import FileResponse
    import os
    
    async with async_session_factory() as db:
        result = await db.execute(
            select(RagDocument).where(
                RagDocument.id == file_id,
                RagDocument.tenant_id == int(tenant_id)
            )
        )
        doc = result.scalar_one_or_none()
        
        if not doc:
            raise HTTPException(status_code=404, detail="文件不存在")
        
        if not os.path.exists(doc.file_path):
            raise HTTPException(status_code=404, detail="物理文件不存在")
        
        return FileResponse(
            path=doc.file_path,
            filename=doc.file_name,
            media_type="application/octet-stream"
        )



# ============== URL爬取API ==============

@router.post("/crawl/url", response_model=URLCrawlResponse)
async def crawl_url(
    request: URLCrawlRequest,
    tenant_id: Optional[str] = Header(None, alias="X-Tenant-ID", description="租户ID")
) -> URLCrawlResponse:
    """
    爬取网页内容
    
    - **url**: 要爬取的URL
    - **extract_images**: 是否提取图片URL
    - **max_depth**: 爬取深度
    """
    import httpx
    from bs4 import BeautifulSoup
    
    try:
        async with httpx.AsyncClient(timeout=30.0, follow_redirects=True) as client:
            response = await client.get(request.url)
            response.raise_for_status()
            
            html = response.text
            soup = BeautifulSoup(html, 'html.parser')
            
            # 提取标题
            title = None
            if soup.title:
                title = soup.title.string
            elif soup.find('h1'):
                title = soup.find('h1').get_text(strip=True)
                
            # 移除脚本和样式
            for script in soup(["script", "style", "nav", "footer", "header"]):
                script.decompose()
            
            # 提取正文内容
            content = soup.get_text(separator='\n', strip=True)
            # 清理多余空白
            content = '\n'.join([line for line in content.split('\n') if line.strip()])
            
            # 提取图片
            images = []
            if request.extract_images:
                for img in soup.find_all('img'):
                    src = img.get('src') or img.get('data-src')
                    if src:
                        if src.startswith('//'):
                            src = 'https:' + src
                        elif src.startswith('/'):
                            # 相对路径转绝对路径
                            from urllib.parse import urljoin
                            src = urljoin(request.url, src)
                        images.append(src)
            
            # 提取元数据
            metadata = {}
            for meta in soup.find_all('meta'):
                name = meta.get('name') or meta.get('property')
                if name and meta.get('content'):
                    metadata[name] = meta.get('content')
            
            return URLCrawlResponse(
                success=True,
                url=request.url,
                title=title,
                content=content[:20000] if content else None,  # 限制内容长度
                metadata=metadata,
                images=images[:20] if images else None  # 限制图片数量
            )
            
    except Exception as e:
        logger.error(f"URL爬取失败: {e}")
        return URLCrawlResponse(
            success=False,
            url=request.url,
            error=str(e)
        )


@router.post("/crawl/url-and-index", response_model=IndexResponse)
async def crawl_and_index_url(
    request: URLCrawlRequest,
    tenant_id: str = Query(..., description="租户ID")
) -> IndexResponse:
    """
    爬取网页并索引到RAG系统
    """
    # 先爬取内容
    crawl_result = await crawl_url(request, tenant_id)
    
    if not crawl_result.success or not crawl_result.content:
        raise HTTPException(status_code=400, detail=crawl_result.error or "爬取失败")
    
    # 获取内容预览（前200字符）
    content_preview = crawl_result.content[:200] + "..." if len(crawl_result.content) > 200 else crawl_result.content
    
    # 索引文档
    doc_id = f"url_{int(datetime.utcnow().timestamp())}"
    document = {
        "id": doc_id,
        "title": crawl_result.title or request.url,
        "description": crawl_result.content,
        "metadata": {
            "source_type": "url_crawl",
            "url": request.url,
            "title": crawl_result.title
        }
    }
    
    success = await enterprise_vector_db.index_patent_document(
        document=document,
        tenant_id=tenant_id
    )
    
    return IndexResponse(
        success=success,
        document_id=doc_id,
        chunk_count=1,
        message=f"标题: {crawl_result.title or '未知'}\n内容预览: {content_preview}\n状态: {'索引成功' if success else '索引失败'}"
    )


# ============== 模型测试API ==============

@router.post("/models/test", response_model=ModelTestResponse)
async def test_embedding_model(
    request: ModelTestRequest
) -> ModelTestResponse:
    """
    测试Embedding模型的连通性
    
    测试逻辑：
    1. 优先使用请求中提供的 api_url 和 api_key
    2. 如果请求没有配置，检查数据库中已保存的自定义模型配置
    3. 如果有 api_url，使用 OpenAI 兼容的 Embedding API 进行测试
    4. 如果没有 api_url，使用本地模型（HuggingFace/Ollama）进行测试
    
    不再根据模型名称推断提供商类型，完全依赖用户配置。
    
    - **model_name**: 模型名称
    - **test_text**: 测试文本（可选）
    - **api_url**: API URL（可选，用于 OpenAI 兼容 API）
    - **api_key**: API Key（可选）
    """
    import time
    import httpx
    
    start_time = time.time()
    model_name = request.model_name
    error_msg = None
    api_url = request.api_url or ""
    api_key = request.api_key or ""
    
    try:
        # 如果请求中没有配置，尝试从数据库获取已保存的配置
        if not api_url:
            from app.ai.rag.config import get_custom_model_config
            saved_config = await get_custom_model_config(model_name)
            if saved_config:
                api_url = saved_config.get("api_url") or ""
                api_key = saved_config.get("api_key") or api_key
        
        test_text = request.test_text or "测试文本"
        
        # 情况1：有 api_url → 使用用户配置的 URL 直接请求
        if api_url and isinstance(api_url, str) and api_url.strip():
            api_url = api_url.strip()
            
            headers = {"Content-Type": "application/json"}
            if api_key:
                headers["Authorization"] = f"Bearer {api_key}"
            
            async with httpx.AsyncClient(timeout=30.0) as client:
                response = await client.post(
                    api_url,
                    headers=headers,
                    json={"input": test_text, "model": model_name}
                )
                
                if response.status_code == 200:
                    result = response.json()
                    data = result.get("data", [{}])
                    embedding_list = data[0].get("embedding", []) if data else []
                    embedding_dim = len(embedding_list)
                    return ModelTestResponse(
                        success=True,
                        model_name=model_name,
                        provider="api",
                        latency_ms=round((time.time() - start_time) * 1000, 2),
                        embedding_dim=embedding_dim,
                        sample_embedding=embedding_list[:10]
                    )
                else:
                    error_msg = f"API返回错误: {response.status_code} - {response.text[:200]}"
        else:
            # 情况2：没有提供 api_url，直接返回错误
            error_msg = "未提供有效的 API URL，请配置自定义模型 API 地址"
        
        return ModelTestResponse(
            success=False,
            model_name=model_name,
            provider="api" if api_url else "local",
            latency_ms=round((time.time() - start_time) * 1000, 2),
            error=error_msg or "未知错误"
        )
        
    except Exception as e:
        logger.error(f"Embedding模型测试失败: {e}")
        return ModelTestResponse(
            success=False,
            model_name=model_name,
            provider="api" if api_url else "local",
            latency_ms=round((time.time() - start_time) * 1000, 2),
            error=str(e)
        )


@router.post("/rerankers/test", response_model=ModelTestResponse)
async def test_reranker_model(
    request: ModelTestRequest
) -> ModelTestResponse:
    """
    测试Reranker模型的连通性
    
    测试逻辑：
    1. 优先使用请求中提供的 api_url 和 api_key
    2. 如果请求没有配置，检查数据库中已保存的自定义模型配置
    3. 如果有 api_url，使用 Rerank API 进行测试（遵循 Cohere/OpenAI 兼容格式）
    4. 如果没有 api_url，使用本地模型（HuggingFace）进行测试
    
    不再根据模型名称推断提供商类型，完全依赖用户配置。
    
    - **model_name**: 模型名称
    - **test_text**: 测试查询文本（可选）
    - **api_url**: API URL（可选，用于 Rerank API）
    - **api_key**: API Key（可选）
    """
    import time
    import httpx
    
    start_time = time.time()
    model_name = request.model_name
    error_msg = None
    api_url = request.api_url or ""
    api_key = request.api_key or ""
    
    # 调试日志
    logger.info(f"Reranker测试请求: model_name={model_name}, request_api_url={request.api_url}")
    
    try:
        # 优先使用请求中传递的配置，如果没有则从数据库获取
        if not api_url:
            from app.ai.rag.config import get_custom_reranker_config
            saved_config = await get_custom_reranker_config(model_name)
            if saved_config and saved_config.get("api_url"):
                api_url = saved_config.get("api_url") or ""
                api_key = saved_config.get("api_key") or api_key
                logger.info(f"从数据库加载配置: api_url={api_url}")
            else:
                logger.warning(f"未找到模型 {model_name} 的配置")
        else:
            logger.info(f"使用请求中传递的API配置")
        
        # 测试数据
        test_query = request.test_text or "什么是人工智能？"
        test_docs = [
            "人工智能是计算机科学的一个分支，旨在创建能够执行需要人类智能的任务的系统。",
            "今天天气很好，适合外出游玩，阳光明媚，微风习习。"
        ]
        
        # 情况1：有 api_url → 使用 Rerank API 测试
        if api_url and isinstance(api_url, str) and api_url.strip():
            api_url = api_url.strip()
            
            headers = {"Content-Type": "application/json"}
            if api_key:
                headers["Authorization"] = f"Bearer {api_key}"
            
            async with httpx.AsyncClient(timeout=30.0) as client:
                response = await client.post(
                    api_url,
                    headers=headers,
                    json={
                        "query": test_query,
                        "documents": test_docs,
                        "top_k": 2,
                        "model": model_name
                    }
                )
                
                if response.status_code == 200:
                    result = response.json()
                    results_list = result.get("results", [])
                    scores = [r.get("relevance_score", 0) for r in results_list[:3]]
                    return ModelTestResponse(
                        success=True,
                        model_name=model_name,
                        provider="api",
                        latency_ms=round((time.time() - start_time) * 1000, 2),
                        embedding_dim=len(results_list),
                        sample_embedding=scores
                    )
                else:
                    error_msg = f"API返回错误: {response.status_code} - {response.text[:200]}"
        else:
            # 情况2：没有提供 api_url，直接返回错误
            error_msg = "未提供有效的 API URL，请配置自定义 Reranker API 地址"
        
        return ModelTestResponse(
            success=False,
            model_name=model_name,
            provider="api" if api_url else "local",
            latency_ms=round((time.time() - start_time) * 1000, 2),
            error=error_msg or "未知错误"
        )
        
    except Exception as e:
        logger.error(f"Reranker模型测试失败: {e}")
        return ModelTestResponse(
            success=False,
            model_name=model_name,
            provider="api" if api_url else "local",
            latency_ms=round((time.time() - start_time) * 1000, 2),
            error=str(e)
        )


# ============== 专利数据源爬取增强 ==============

@router.post("/patents/crawl", response_model=PatentSearchAndIndexResponse)
async def crawl_patents(
    request: PatentSearchRequest,
    tenant_id: str = Query(..., description="租户ID")
) -> PatentSearchAndIndexResponse:
    """
    爬取公开专利数据并索引
    
    支持的数据源:
    - **dawei**: 大为专利数据库 (中国)
    - **cnipa**: CNIPA官方数据库 (中国)
    - **uspto**: USPTO PatentsView (美国)
    - **wipo**: WIPO PATENTSCOPE (国际)
    - **epo**: 欧洲专利局
    """
    from app.ai.patent_database_api import patent_aggregator
    
    # 收集每个数据源的错误信息
    source_errors = {}
    source_results = {}
    
    try:
        sources_arg: Optional[List[str]] = request.sources
        results = await patent_aggregator.search_all(
            query=request.query,
            sources=sources_arg,
            max_results_per_source=request.max_results
        )
        
        # 收集所有专利
        all_patents = []
        for source, patents in results.items():
            if patents:
                all_patents.extend(patents)
                source_results[source] = len(patents)
            else:
                # 检查是否有错误
                source_errors[source] = "未找到专利或API不可用"
        
        # 索引到向量数据库
        indexed_count = 0
        index_details = {}
        if all_patents:
            docs = [p.to_dict() for p in all_patents]
            index_results = await enterprise_vector_db.index_patent_batch(
                documents=docs,
                tenant_id=tenant_id
            )
            indexed_count = sum(1 for v in index_results.values() if v)
            index_details = {k: v for k, v in index_results.items()}
        
        # 构建响应消息
        message_parts = []
        if source_results:
            message_parts.append(f"各数据源找到: {', '.join([f'{k}:{v}条' for k,v in source_results.items()])}")
        if source_errors:
            message_parts.append(f"警告: {', '.join([f'{k}:{v}' for k,v in source_errors.items()])}")
        if indexed_count > 0:
            message_parts.append(f"成功索引: {indexed_count}条")
        
        return PatentSearchAndIndexResponse(
            total_found=len(all_patents),
            indexed=indexed_count,
            sources=source_results,
            message="\n".join(message_parts) if message_parts else "爬取完成"
        )
        
    except Exception as e:
        logger.error(f"专利爬取失败: {e}")
        raise HTTPException(status_code=500, detail=str(e))

# ========== 搜索历史管理 ==========

@router.get("/search-history")
async def get_search_history(
    request: Request,
    tenant_id: str = Query(..., description="租户ID"),
    limit: int = Query(50, ge=1, le=200, description="返回数量限制")
):
    """获取搜索历史列表"""
    try:
        from sqlalchemy import select, func
        from app.database.engine import async_session_factory
        from app.database.models import SearchHistory, User
        
        async with async_session_factory() as db:
            # 获取总数量
            count_query = select(func.count()).select_from(SearchHistory).where(
                SearchHistory.tenant_id == int(tenant_id)
            )
            count_result = await db.execute(count_query)
            total = count_result.scalar()
            
            # 获取历史记录，按时间倒序
            query = select(SearchHistory).where(
                SearchHistory.tenant_id == int(tenant_id)
            ).order_by(SearchHistory.timestamp.desc()).limit(limit)
            
            result = await db.execute(query)
            history_records = result.scalars().all()
            
            # 转换格式，添加用户名
            history = []
            for record in history_records:
                item = {
                    "id": str(record.id),
                    "query": record.query,
                    "search_type": record.search_type,
                    "result_count": record.result_count,
                    "timestamp": record.timestamp.isoformat(),
                    "user_name": record.user.username if record.user else None
                }
                history.append(item)
            
            return {
                "history": history,
                "total": total
            }
    except Exception as e:
        logger.error(f"获取搜索历史失败: {str(e)}")
        raise HTTPException(status_code=500, detail=f"获取搜索历史失败: {str(e)}")

@router.delete("/search-history/{history_id}")
async def delete_search_history(
    request: Request,
    history_id: int = Path(..., description="历史记录ID"),
    tenant_id: str = Query(..., description="租户ID")
):
    """删除单条搜索历史"""
    try:
        from sqlalchemy import delete
        from app.database.engine import async_session_factory
        from app.database.models import SearchHistory
        
        async with async_session_factory() as db:
            stmt = delete(SearchHistory).where(
                SearchHistory.id == history_id,
                SearchHistory.tenant_id == int(tenant_id)
            )
            result = await db.execute(stmt)
            await db.commit()
            
            if result.rowcount == 0:
                raise HTTPException(status_code=404, detail="未找到历史记录")
            
            return {"success": True, "message": "已删除"}
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"删除搜索历史失败: {str(e)}")
        raise HTTPException(status_code=500, detail=f"删除搜索历史失败: {str(e)}")

@router.delete("/search-history")
async def clear_search_history(
    request: Request,
    tenant_id: str = Query(..., description="租户ID")
):
    """清空搜索历史"""
    try:
        from sqlalchemy import delete
        from app.database.engine import async_session_factory
        from app.database.models import SearchHistory
        
        async with async_session_factory() as db:
            stmt = delete(SearchHistory).where(
                SearchHistory.tenant_id == int(tenant_id)
            )
            await db.execute(stmt)
            await db.commit()
            
            return {"success": True, "message": "已清空历史"}
    except Exception as e:
        logger.error(f"清空搜索历史失败: {str(e)}")
        raise HTTPException(status_code=500, detail=f"清空搜索历史失败: {str(e)}")